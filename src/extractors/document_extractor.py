import os
from typing import Optional
    
def extract_pdf(file_path:str) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text.strip()

def extract_docx(file_path: str) -> str:
    import docx
    doc = docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()
    
def extract_epub(file_path: str) -> str:
    import ebooklib
    from ebooklib import epub

    book = epub.read_epub(file_path)
    text = ""
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            text += item.content.decode("utf-8")
    return text.strip()

def extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text.strip()